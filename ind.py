from pptx import Presentation
from pptx.dml.color import RGBColor
prs = Presentation("test.pptx")
from pptx.enum.dml import MSO_THEME_COLOR
from utilscripts import gradientsetter
# title_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(title_slide_layout)
# title = slide.shapes.title
# subtitle = slide.placeholders[1]

# title.text = "Hello, World!"
# subtitle.text = "python-pptx was here!"


for slide in prs.slides:
    print(slide.slide_id)
    print(slide.background)
    
    background = slide.background
    fill = background.fill
    fila = gradientsetter(fill, 120, RGBColor(255, 196, 0), RGBColor(250, 0, 158))
    fill = fila
    print(fill.gradient_stops[0].color.theme_color)

    title = slide.shapes.title.text
    slide.notes_slide.notes_text_frame.text = "This is the note for slide " + str(slide.slide_id) + " with title " + title
    for shape in slide.shapes:
        if shape.has_text_frame:
            print(shape.text)
            shape.text = shape.text + "!!@@@"
        else:
            print("No text frame")

prs.save('test.pptx')